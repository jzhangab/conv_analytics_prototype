"""
Generate a PowerPoint deck explaining the Clinical Analytics Chatbot architecture,
data flows, and agentic AI workflows.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ───────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1A, 0x29, 0x4A)   # slide background / title bg
BLUE      = RGBColor(0x4A, 0x90, 0xD9)   # user-facing / accent
PURPLE    = RGBColor(0x7B, 0x68, 0xEE)   # orchestrator
AMBER     = RGBColor(0xE8, 0xA8, 0x38)   # LLM
GREEN     = RGBColor(0x3C, 0xB3, 0x71)   # agents
RED       = RGBColor(0xCD, 0x5C, 0x5C)   # data
SLATE     = RGBColor(0x70, 0x80, 0x90)   # output
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHTGRAY = RGBColor(0xF0, 0xF2, 0xF5)
DARKGRAY  = RGBColor(0x33, 0x33, 0x33)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


# ── Low-level helpers ────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_pt=1.0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=14, bold=False, color=DARKGRAY,
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_label_box(slide, text, left, top, width, height,
                  bg=LIGHTGRAY, fg=DARKGRAY, font_size=11, bold=False,
                  line_rgb=None, valign_center=True, align=PP_ALIGN.CENTER):
    """Filled rounded rectangle with centered text."""
    rect = slide.shapes.add_shape(
        5,  # rounded rectangle
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    rect.adjustments[0] = 0.05
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    if line_rgb:
        rect.line.color.rgb = line_rgb
        rect.line.width = Pt(1.2)
    else:
        rect.line.fill.background()

    tf = rect.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = fg
    return rect


def slide_bg(slide, color=NAVY):
    bg = slide.shapes.add_shape(
        1, Emu(0), Emu(0), SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.zorder = 0


def section_bar(slide, title, left=0.2, top=0.3, width=12.9, height=0.38,
                color=BLUE):
    rect = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    tf = rect.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = title
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE


def slide_title(slide, title, subtitle=None):
    add_text_box(slide, title, 0.4, 0.12, 12.5, 0.55,
                 font_size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.4, 0.68, 12.5, 0.35,
                     font_size=12, bold=False, color=RGBColor(0xBB, 0xCC, 0xEE),
                     align=PP_ALIGN.LEFT)


def add_bullet_box(slide, items, left, top, width, height,
                   font_size=11, color=WHITE, indent=0.2, spacing=0.32):
    """items: list of (indent_level, text). level 0 = bullet, level 1 = sub-bullet."""
    y = top
    for level, text in items:
        bullet = "•" if level == 0 else "  ◦"
        add_text_box(slide, f"{bullet}  {text}", left + level*indent, y,
                     width - level*indent, spacing,
                     font_size=font_size, color=color)
        y += spacing


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Title slide."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    slide_bg(slide, NAVY)

    # Accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(3.2), SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    add_text_box(slide, "Clinical Analytics Chatbot",
                 0.5, 3.3, 12.3, 0.7,
                 font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Architecture · Data Flows · Agentic AI Workflows",
                 0.5, 4.05, 12.3, 0.45,
                 font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Dataiku DSS  ·  Flask Webapp  ·  LLM Mesh (Azure OpenAI)",
                 0.5, 5.0, 12.3, 0.35,
                 font_size=12, color=RGBColor(0xBB, 0xCC, 0xEE), align=PP_ALIGN.CENTER)
    add_text_box(slide, "Drug R&D · Clinical Operations",
                 0.5, 5.45, 12.3, 0.3,
                 font_size=11, color=RGBColor(0x99, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


def slide_02_overview(prs):
    """Project overview."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    slide_bg(slide, NAVY)
    slide_title(slide, "Project Overview",
                "An agentic AI chatbot for drug R&D clinical analytics, deployed on Dataiku DSS")

    section_bar(slide, "What It Does", top=1.1, color=BLUE)
    items = [
        (0, "Natural-language interface for clinical operations teams"),
        (0, "Seven specialist AI agents cover the full trial planning workflow"),
        (0, "Orchestrator drives intent classification → parameter gathering → confirmation → execution"),
        (0, "Results displayed in chat; optionally exported to Dataiku datasets on user request"),
        (0, "Grounded data-reasoning: ask follow-up questions across all prior results in a session"),
    ]
    add_bullet_box(slide, items, 0.4, 1.55, 12.5, 3.0, font_size=12)

    section_bar(slide, "Seven Analytical Skills", top=4.05, color=GREEN)
    skills = [
        ("CRO Site Profiling",      BLUE),
        ("Trial Benchmarking",      PURPLE),
        ("Drug Reimbursement",      AMBER),
        ("Enrollment Forecasting",  GREEN),
        ("Protocol Analysis",       RED),
        ("Country Ranking",         SLATE),
        ("Enrollment Reforecasting",RGBColor(0x20, 0x8D, 0xA0)),
    ]
    for i, (name, color) in enumerate(skills):
        col = i % 4
        row = i // 4
        add_label_box(slide, name,
                      0.35 + col * 3.2, 4.55 + row * 0.82,
                      3.0, 0.65,
                      bg=color, fg=WHITE, font_size=10, bold=True)


def slide_03_architecture(prs):
    """High-level architecture layers."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    slide_bg(slide, NAVY)
    slide_title(slide, "System Architecture",
                "Seven layers from browser to data — all running inside Dataiku DSS")

    layers = [
        (BLUE,   "① User Interface",      "Browser → Flask Webapp (GET /, POST /chat, /upload, /confirm, /export)"),
        (PURPLE, "② Orchestrator + FSM",  "orchestrator.py — drives the state machine, coordinates all components"),
        (AMBER,  "③ LLM Client",          "llm_client.py → Dataiku LLM Mesh → Azure OpenAI (max 16,384 tokens)"),
        (AMBER,  "   Web Search",         "web_search.py — supplementary context for benchmarking & reasoning"),
        (GREEN,  "④ Skill Agents (×7)",   "router.py → agents/ — each agent owns its data source + LLM prompts"),
        (RED,    "⑤ Agent Outputs",       "AgentResult → SkillResult stored in ConversationState"),
        (SLATE,  "⑥ Response / Export",   "_serialize_response() → JSON to UI; POST /export → Dataiku Dataset"),
        (RGBColor(0x88,0x88,0x88),
                 "⑦ Configuration",       "llm_config.yaml · skills_config.yaml · shared param inheritance"),
    ]
    for i, (color, name, desc) in enumerate(layers):
        y = 1.1 + i * 0.73
        # colour strip
        strip = slide.shapes.add_shape(1, Inches(0.3), Inches(y),
                                       Inches(0.18), Inches(0.55))
        strip.fill.solid()
        strip.fill.fore_color.rgb = color
        strip.line.fill.background()
        # name
        add_text_box(slide, name, 0.6, y + 0.04, 3.2, 0.52,
                     font_size=11, bold=True, color=WHITE)
        # description
        add_text_box(slide, desc, 3.85, y + 0.04, 9.1, 0.52,
                     font_size=10.5, color=RGBColor(0xCC, 0xDD, 0xFF))


def slide_04_fsm(prs):
    """Orchestrator FSM states and transitions."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    slide_bg(slide, NAVY)
    slide_title(slide, "Orchestrator — Finite State Machine",
                "Every user message is routed through a deterministic state machine before any LLM call")

    # FSM state boxes
    states = ["IDLE", "PARAMETER\nGATHERING", "CONFIRMATION\nPENDING",
              "SKILL\nEXECUTION", "ANALYSIS\nPLANNING"]
    colors = [BLUE, PURPLE, AMBER, GREEN, RED]
    xs     = [0.35, 2.65, 5.05, 7.45, 9.85]

    section_bar(slide, "FSM States", top=1.1, width=12.5, color=PURPLE)
    for i, (state, color) in enumerate(zip(states, colors)):
        add_label_box(slide, state, xs[i], 1.6, 2.1, 1.0,
                      bg=color, fg=WHITE, font_size=11, bold=True)
        if i < len(states) - 1:
            add_text_box(slide, "→", xs[i]+2.12, 1.9, 0.4, 0.45,
                         font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Back-to-IDLE note
    add_text_box(slide, "⟳ SKILL EXECUTION → IDLE  (after result or cancel)",
                 0.35, 2.7, 9.5, 0.35, font_size=10,
                 color=RGBColor(0xBB, 0xCC, 0xFF), align=PP_ALIGN.CENTER)

    section_bar(slide, "Per-Request Flow", top=3.15, width=12.5, color=PURPLE)
    steps = [
        "1  Retrieve or create ConversationState for session_id",
        "2  Add user message to history",
        "3  Route by current FSM state:",
        "     IDLE → classify_intent() via LLM  →  set active_skill",
        "     PARAMETER_GATHERING → extract_parameters() via LLM  →  merge into state",
        "     CONFIRMATION_PENDING → parse yes / no / edit  →  execute or retry",
        "     SKILL_EXECUTION → router.get_agent(skill_id).run(params, state)",
        "     ANALYSIS_PLANNING → confirm plan  →  _handle_reasoning() with web search",
        "4  Build JSON response  →  return to Flask  →  browser",
    ]
    for i, step in enumerate(steps):
        indent = 0.7 if step.startswith(" ") else 0.4
        add_text_box(slide, step.strip(), indent, 3.6 + i*0.36, 12.1, 0.34,
                     font_size=10.5,
                     color=WHITE if not step.startswith(" ") else RGBColor(0xBB,0xCC,0xFF))

    # Shared-param inheritance note
    add_text_box(slide,
                 "Shared parameter inheritance: indication · age_group · phase "
                 "auto-propagated across skills within a session",
                 0.4, 6.95, 12.3, 0.35, font_size=9.5,
                 color=RGBColor(0x99, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


def slide_05_llm(prs):
    """LLM Client & Dataiku LLM Mesh."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    slide_bg(slide, NAVY)
    slide_title(slide, "LLM Client & Dataiku LLM Mesh",
                "All model calls are routed through Dataiku's LLM abstraction layer")

    section_bar(slide, "LLMClient (llm_client.py)", top=1.1, width=12.5, color=AMBER)

    methods = [
        ("complete()", "messages list → text response"),
        ("complete_json()", "complete() + strip fences + json.loads"),
        ("_repair_json()", "close open strings / arrays / objects (truncation recovery)"),
    ]
    for i, (method, desc) in enumerate(methods):
        add_label_box(slide, method, 0.4 + i*4.15, 1.6, 3.8, 0.55,
                      bg=RGBColor(0x4A,0x3A,0x10), fg=AMBER, font_size=11, bold=True,
                      line_rgb=AMBER)
        add_text_box(slide, desc, 0.4 + i*4.15, 2.2, 3.8, 0.4,
                     font_size=10, color=RGBColor(0xCC,0xDD,0xFF))

    section_bar(slide, "Dataiku LLM Mesh Call Chain", top=2.75, width=12.5, color=AMBER)
    chain = [
        "project.get_llm(connection_id)",
        "llm.new_completion()",
        "completion.add_message() for each turn",
        "completion.with_max_output_tokens(16384)",
        "completion.execute()  →  resp.text",
    ]
    for i, step in enumerate(chain):
        col = i % 3
        row = i // 3
        add_label_box(slide, step,
                      0.4 + col * 4.15, 3.25 + row * 0.68, 3.8, 0.55,
                      bg=RGBColor(0x3A, 0x2A, 0x00), fg=WHITE, font_size=10,
                      line_rgb=AMBER)

    section_bar(slide, "Prompt Templates (prompt_templates.py)", top=4.7, width=12.5, color=AMBER)
    templates = [
        "INTENT_CLASSIFIER", "PARAMETER_EXTRACTOR", "SITE_MATCHING",
        "TRIAL_BENCHMARKING", "DRUG_REIMBURSEMENT", "ENROLLMENT_PARAMS",
        "ENROLLMENT_NARRATIVE", "ANALYSIS_PLAN", "DATA_REASONING", "GENERAL_KNOWLEDGE",
    ]
    for i, t in enumerate(templates):
        col = i % 5
        row = i // 5
        add_label_box(slide, t, 0.35 + col*2.52, 5.2 + row*0.6, 2.38, 0.48,
                      bg=RGBColor(0x28, 0x20, 0x00), fg=AMBER, font_size=9)

    add_text_box(slide,
                 "Web search (web_search.py) provides supplementary context to "
                 "Trial Benchmarking, Drug Reimbursement, Protocol Analysis, and Data Reasoning calls",
                 0.4, 6.95, 12.3, 0.35, font_size=9.5,
                 color=RGBColor(0x99, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


def slide_06_agents(prs):
    """Skill agents deep-dive."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    slide_bg(slide, NAVY)
    slide_title(slide, "Seven Skill Agents",
                "Each agent owns its data source, LLM prompts, and output schema")

    section_bar(slide, "Agent Registry  (router.py → agents/)", top=1.1, width=12.5, color=GREEN)

    agents = [
        (GREEN,  "CRO Site\nProfiling",
         "site_list_merger_agent.py",
         "2-step Jaro-Winkler match of uploaded CRO site list vs CTMS master DB\n"
         "Returns: confidence score, PI, avg enrolled, avg months diff per site"),
        (PURPLE, "Trial\nBenchmarking",
         "trial_benchmarking_agent.py",
         "Queries Citeline DB (indication + phase + age_group) + web search\n"
         "Returns: key metrics, patterns, challenges, narrative summary"),
        (AMBER,  "Drug\nReimbursement",
         "drug_reimbursement_agent.py",
         "HTA/payer landscape for user-provided countries (no defaults)\n"
         "Returns: country-level likelihood, requirements, risk factors"),
        (RED,    "Enrollment\nForecasting",
         "enrollment_forecasting_agent.py",
         "Stage 1: LLM estimates domain params per scenario\n"
         "Stage 2: deterministic logistic ramp → 3 Bokeh curves (pessimistic / moderate / optimistic)"),
        (BLUE,   "Protocol\nAnalysis",
         "protocol_analysis_agent.py",
         "Reviews uploaded protocol PDF or text for study design risks\n"
         "Returns: structured design review, operational risks, recommendations"),
        (SLATE,  "Country\nRanking",
         "country_ranking_agent.py",
         "Ranks countries by clinical trial experience for a given indication\n"
         "Returns: ranked table with experience score and notes"),
        (RGBColor(0x20,0x8D,0xA0), "Enrollment\nReforecasting",
         "reforecasting_agent.py",
         "Updates enrollment forecast using actual site activation data\n"
         "Returns: revised curves vs original, variance analysis"),
    ]

    for i, (color, name, module, desc) in enumerate(agents):
        col = i % 2
        row = i // 2
        x = 0.35 + col * 6.45
        y = 1.6 + row * 1.35
        add_label_box(slide, name, x, y, 1.55, 1.15,
                      bg=color, fg=WHITE, font_size=10, bold=True)
        add_text_box(slide, module, x + 1.65, y + 0.0, 4.6, 0.28,
                     font_size=9, bold=True, color=color)
        add_text_box(slide, desc, x + 1.65, y + 0.28, 4.6, 0.85,
                     font_size=9.0, color=RGBColor(0xCC, 0xDD, 0xFF))

    add_text_box(slide,
                 "All agents extend BaseAgent (base_agent.py) and return AgentResult "
                 "(success, text_response, table_data, table_columns, chart_json)",
                 0.4, 7.08, 12.3, 0.32, font_size=9.5,
                 color=RGBColor(0x99, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


def slide_07_data_flows(prs):
    """Data flows — inputs, processing, outputs."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    slide_bg(slide, NAVY)
    slide_title(slide, "Data Flows",
                "How data enters, is transformed, and exits the system")

    section_bar(slide, "Inputs", top=1.1, width=4.0, color=BLUE)
    inputs = [
        (0, "User text messages (POST /chat)"),
        (0, "site_file upload — CSV/Excel CRO site list"),
        (0, "protocol_file upload — PDF or text protocol document"),
        (0, "POST /confirm — yes / no / edit_params JSON"),
        (0, "POST /export — result_id + dataset_name"),
    ]
    add_bullet_box(slide, inputs, 0.35, 1.55, 4.0, 2.5, font_size=10.5)

    section_bar(slide, "Internal Data Stores", top=1.1, left=4.5, width=4.0, color=PURPLE)
    stores = [
        (0, "ConversationState (in-memory per session)"),
        (1, "messages history (last 10 turns to LLM)"),
        (1, "collected_parameters per skill"),
        (1, "uploaded_files (parsed DataFrames)"),
        (1, "prior_results list of SkillResult"),
        (0, "CTMS_SITES.csv — sponsor site master"),
        (0, "Citeline DB — trial benchmarks"),
        (0, "REFORECAST dataset — actual activation data"),
    ]
    add_bullet_box(slide, stores, 4.55, 1.55, 4.0, 3.0, font_size=10.5)

    section_bar(slide, "Outputs", top=1.1, left=8.7, width=4.2, color=GREEN)
    outputs = [
        (0, "JSON response to browser"),
        (1, "message (Markdown narrative)"),
        (1, "table_data + table_columns"),
        (1, "chart_json (Bokeh, forecasting)"),
        (1, "result_id + fsm_state"),
        (0, "Dataiku Dataset (on explicit /export)"),
        (0, "Confirmation dialog prompt to UI"),
    ]
    add_bullet_box(slide, outputs, 8.75, 1.55, 4.1, 3.0, font_size=10.5)

    section_bar(slide, "Data Transformation Pipeline (per skill)", top=4.7, width=12.5, color=RED)

    pipeline = [
        ("Raw Input\n(file / text)", BLUE),
        ("Parameter\nExtraction (LLM)", PURPLE),
        ("Agent Logic\n+ Data Query", GREEN),
        ("LLM Synthesis\n+ Narrative", AMBER),
        ("AgentResult\n→ SkillResult", RED),
        ("JSON Response\n→ Browser", SLATE),
    ]
    for i, (label, color) in enumerate(pipeline):
        add_label_box(slide, label, 0.35 + i*2.14, 5.2, 1.9, 0.85,
                      bg=color, fg=WHITE, font_size=10, bold=True)
        if i < len(pipeline) - 1:
            add_text_box(slide, "→", 0.35 + i*2.14 + 1.92, 5.45, 0.2, 0.38,
                         font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 "File uploads parsed to DataFrames in memory; "
                 "results persist in session for follow-up reasoning; "
                 "Dataiku dataset export only on explicit user request",
                 0.4, 6.2, 12.3, 0.45, font_size=9.5,
                 color=RGBColor(0x99,0xAA,0xCC), align=PP_ALIGN.CENTER)


def slide_08_agentic_workflow(prs):
    """Agentic AI workflow — the end-to-end loop."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    slide_bg(slide, NAVY)
    slide_title(slide, "Agentic AI Workflow",
                "How the system orchestrates multiple LLM calls to complete a user request")

    section_bar(slide, "Standard Skill Execution (e.g. Enrollment Forecasting)",
                top=1.1, width=12.5, color=GREEN)

    standard_steps = [
        ("User", BLUE,   "Types: \"Forecast enrollment for NSCLC adults Phase 3\""),
        ("LLM",  AMBER,  "classify_intent() → skill_id = 'enrollment_forecasting'"),
        ("LLM",  AMBER,  "extract_parameters() → {indication, age_group, phase, n_sites, …}"),
        ("Orch", PURPLE, "build_confirmation_prompt() → show parameter summary to user"),
        ("User", BLUE,   "Confirms → POST /confirm {confirmed: true}"),
        ("Agent",GREEN,  "EnrollmentForecastingAgent.run() — Stage 1: LLM estimates domain params"),
        ("Agent",GREEN,  "Stage 2: deterministic logistic math → 3 scenario curves (pessimistic/moderate/optimistic)"),
        ("LLM",  AMBER,  "ENROLLMENT_NARRATIVE prompt → narrative text with peak sites, months-to-target"),
        ("UI",   BLUE,   "Bokeh chart + narrative returned as JSON → rendered in browser"),
    ]
    for i, (actor, color, desc) in enumerate(standard_steps):
        add_label_box(slide, actor, 0.35, 1.6 + i*0.52, 0.7, 0.42,
                      bg=color, fg=WHITE, font_size=9, bold=True)
        add_text_box(slide, desc, 1.15, 1.63 + i*0.52, 11.5, 0.4,
                     font_size=10, color=WHITE)

    section_bar(slide, "Data Reasoning Workflow (follow-up question on prior results)",
                top=6.42, width=12.5, color=RED)
    reasoning = ("Prior SkillResults in session  →  _generate_plan() via LLM  →  "
                 "User confirms / revises plan  →  _handle_reasoning() with DATA_REASONING prompt "
                 "+ web search  →  grounded analytical answer")
    add_text_box(slide, reasoning, 0.4, 6.88, 12.3, 0.45,
                 font_size=10, color=WHITE)


def slide_09_file_upload(prs):
    """File upload + confirmation flows."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    slide_bg(slide, NAVY)
    slide_title(slide, "File Upload & Confirmation Flows",
                "Two dedicated upload types; every skill execution requires explicit user confirmation")

    section_bar(slide, "File Upload Flow  (POST /upload)", top=1.1, width=6.0, color=BLUE)
    upload_steps = [
        "1  Browser sends multipart/form-data with file_key",
        "2  file_key = 'site_file'  →  parse_uploaded_file()",
        "     Returns: filename, rows, columns (DataFrame in memory)",
        "3  file_key = 'protocol_file'  →  parse_protocol_file()",
        "     Returns: filename, full_text / pages (PDF or text)",
        "4  ConversationState.uploaded_files[file_key] = file_info",
        "5  FSM transitions to PARAMETER_GATHERING for the relevant skill",
        "6  Assistant replies: 'File uploaded: X rows, columns: …'",
    ]
    for i, step in enumerate(upload_steps):
        indent = 0.7 if step.startswith(" ") else 0.4
        add_text_box(slide, step.strip(), indent, 1.55 + i*0.52, 5.7, 0.45,
                     font_size=10.5,
                     color=WHITE if not step.startswith(" ") else RGBColor(0xBB,0xCC,0xFF))

    section_bar(slide, "Confirmation Flow  (POST /confirm)", top=1.1, left=6.4, width=6.5, color=PURPLE)
    confirm_steps = [
        "1  Orchestrator builds human-readable parameter summary",
        "2  FSM state = CONFIRMATION_PENDING",
        "3  Browser shows confirm dialog with param snapshot",
        "4a  confirmed=true  →  _execute_skill(state)",
        "4b  confirmed=false  →  cancel, reset to IDLE",
        "4c  edit_params={…}  →  merge new params, re-check completeness",
        "5  Agent.run(params, state)  →  AgentResult",
        "6  SkillResult stored; FSM → IDLE; JSON response to UI",
    ]
    for i, step in enumerate(confirm_steps):
        indent = 0.7 if step.startswith(" ") else 0.0
        add_text_box(slide, step.strip(), 6.6 + indent*0.3, 1.55 + i*0.52, 6.1, 0.45,
                     font_size=10.5,
                     color=WHITE if not step.startswith(" ") else RGBColor(0xBB,0xCC,0xFF))

    section_bar(slide, "Key Design Decisions", top=5.9, width=12.5,
                color=RGBColor(0x88,0x88,0x88))
    decisions = [
        "File upload — CSV/Excel only, parsed to in-memory DataFrame (no Dataiku dataset reference)",
        "Countries for Drug Reimbursement — always user-provided; no default country list",
        "Enrollment forecasting — always produces three curves: pessimistic, moderate, optimistic",
        "Result persistence — written to Dataiku dataset ONLY on explicit POST /export; displayed in chat first",
        "Access control — none; session isolation via UUID session_id",
    ]
    for i, d in enumerate(decisions):
        add_text_box(slide, f"•  {d}", 0.4, 6.35 + i*0.22, 12.3, 0.22,
                     font_size=9.5, color=RGBColor(0xCC,0xDD,0xFF))


def slide_10_deployment(prs):
    """Deployment & configuration."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    slide_bg(slide, NAVY)
    slide_title(slide, "Deployment & Configuration",
                "Runs as a Dataiku DSS Code Webapp (Flask backend)")

    section_bar(slide, "Dataiku Deployment", top=1.1, width=6.0, color=BLUE)
    deploy = [
        (0, "Backend type: Flask  (webapp.py entry point)"),
        (0, "Lazy initialization — deferred until first request"),
        (1, "Avoids import-time failures surfacing as silent 500s"),
        (0, "SessionStore: in-memory dict, 30-minute TTL"),
        (0, "Frontend: Jinja2 templates + static JS/CSS"),
        (0, "Secret key via FLASK_SECRET_KEY env var"),
        (0, "Health check: GET /healthz → {status: ok}"),
    ]
    add_bullet_box(slide, deploy, 0.35, 1.55, 5.8, 3.5, font_size=11)

    section_bar(slide, "Configuration Files", top=1.1, left=6.4, width=6.5, color=AMBER)
    configs = [
        (0, "llm_config.yaml"),
        (1, "connection_id: YOUR_LLM_CONNECTION_ID"),
        (1, "max_tokens: 16384"),
        (1, "context_window_turns: 10"),
        (1, "temperatures: agents / extraction / narrative"),
        (0, "skills_config.yaml"),
        (1, "7 skill schemas with required / optional params"),
        (1, "param types: string, list, int, file"),
        (1, "choices lists for controlled vocabularies"),
    ]
    add_bullet_box(slide, configs, 6.5, 1.55, 6.2, 3.5, font_size=11)

    section_bar(slide, "File Structure", top=4.9, width=12.5,
                color=RGBColor(0x55,0x55,0x77))
    structure = (
        "webapp.py  ·  config/llm_config.yaml  ·  config/skills_config.yaml  ·  "
        "backend/orchestrator/{orchestrator, router, intent_classifier, parameter_extractor, confirmation_manager}.py  ·  "
        "backend/agents/{base_agent, site_list_merger, trial_benchmarking, drug_reimbursement, "
        "enrollment_forecasting, protocol_analysis, country_ranking, reforecasting}_agent.py  ·  "
        "backend/llm/{llm_client, prompt_templates, response_parser, web_search}.py  ·  "
        "backend/state/{conversation_state, session_store, parameter_schema}.py  ·  "
        "frontend/templates/index.html  ·  frontend/static/{css,js}/"
    )
    add_text_box(slide, structure, 0.4, 5.35, 12.3, 1.8,
                 font_size=9, color=RGBColor(0xBB,0xCC,0xFF))


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()
    slide_01_title(prs)
    slide_02_overview(prs)
    slide_03_architecture(prs)
    slide_04_fsm(prs)
    slide_05_llm(prs)
    slide_06_agents(prs)
    slide_07_data_flows(prs)
    slide_08_agentic_workflow(prs)
    slide_09_file_upload(prs)
    slide_10_deployment(prs)

    out = "/Users/jeremyzhang/conv_analytics_prototype/conv_analytics_architecture.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
